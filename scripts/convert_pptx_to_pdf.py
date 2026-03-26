#!/usr/bin/env python3
"""Convert PPTX to PDF with tool-first strategy and text-PDF fallback."""

from __future__ import annotations

import argparse
import json
import shutil
import subprocess
from datetime import datetime, timezone
from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Convert a PPTX deck into PDF.")
    p.add_argument("--input-pptx", type=Path, required=True)
    p.add_argument("--output-pdf", type=Path, required=True)
    p.add_argument(
        "--report-json",
        type=Path,
        default=None,
        help="Optional JSON report path for conversion metadata.",
    )
    p.add_argument(
        "--strategy",
        choices=["auto", "soffice", "keynote", "fallback"],
        default="auto",
        help="auto tries soffice/libreoffice, then keynote, then text fallback.",
    )
    p.add_argument(
        "--require-visual",
        action="store_true",
        help="Fail if conversion cannot be done by visual renderers (soffice/keynote).",
    )
    return p.parse_args()


def run(cmd: list[str], timeout_sec: int = 30) -> subprocess.CompletedProcess[str]:
    try:
        return subprocess.run(cmd, text=True, capture_output=True, check=False, timeout=timeout_sec)
    except subprocess.TimeoutExpired as e:
        return subprocess.CompletedProcess(
            cmd,
            124,
            (e.stdout or "") if isinstance(e.stdout, str) else "",
            ((e.stderr or "") if isinstance(e.stderr, str) else "") + f"\nTimeout after {timeout_sec}s",
        )


def _tool_path(name: str) -> str | None:
    return shutil.which(name)


def try_soffice(input_pptx: Path, output_pdf: Path) -> tuple[bool, str]:
    bin_path = _tool_path("soffice") or _tool_path("libreoffice")
    if not bin_path:
        return False, "soffice/libreoffice not found"

    out_dir = output_pdf.parent
    expected = out_dir / f"{input_pptx.stem}.pdf"
    if expected.exists():
        expected.unlink()

    cp = run([bin_path, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(input_pptx)], timeout_sec=90)
    if cp.returncode != 0:
        return False, f"soffice failed: {cp.stderr.strip() or cp.stdout.strip() or 'unknown error'}"
    if not expected.exists():
        return False, "soffice returned success but output PDF missing"
    if expected.resolve() != output_pdf.resolve():
        if output_pdf.exists():
            output_pdf.unlink()
        expected.rename(output_pdf)
    return True, "soffice"


def try_keynote(input_pptx: Path, output_pdf: Path) -> tuple[bool, str]:
    if not _tool_path("osascript"):
        return False, "osascript not found"

    # Keep this script compact and defensive. Some environments do not allow
    # opening external files from Keynote automation; failures are expected.
    script = f"""
set inputFile to POSIX file "{input_pptx}"
set outputFile to POSIX file "{output_pdf}"
tell application "Keynote"
    activate
    try
        set theDoc to open inputFile
    on error
        return "keynote_open_failed"
    end try
    delay 2
    try
        export theDoc to outputFile as PDF
        close theDoc saving no
    on error
        return "keynote_export_failed"
    end try
end tell
return "ok"
"""
    cp = run(["osascript", "-e", script], timeout_sec=15)
    if cp.returncode != 0:
        return False, f"keynote osascript failed: {cp.stderr.strip() or cp.stdout.strip() or 'unknown error'}"
    marker = (cp.stdout or "").strip()
    if marker != "ok":
        return False, f"keynote failed marker: {marker}"
    if not output_pdf.exists():
        return False, "keynote returned ok but output PDF missing"
    return True, "keynote"


def fallback_text_pdf(input_pptx: Path, output_pdf: Path) -> tuple[bool, str]:
    prs = Presentation(str(input_pptx))
    output_pdf.parent.mkdir(parents=True, exist_ok=True)

    with PdfPages(output_pdf) as pdf:
        for i, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shp in slide.shapes:
                if hasattr(shp, "text"):
                    t = (shp.text or "").strip()
                    if t:
                        texts.append(t.replace("\r", "\n"))
            body = "\n\n".join(texts) if texts else "(no text found on this slide)"

            fig = plt.figure(figsize=(11.69, 8.27), dpi=150)
            fig.patch.set_facecolor("white")
            ax = fig.add_axes([0.05, 0.06, 0.90, 0.90])
            ax.axis("off")
            ax.text(0.0, 0.98, "Text-based PDF fallback (slide visuals not rendered)", fontsize=10, color="#6b6b6b", va="top")
            ax.text(0.0, 0.93, f"Slide {i}/{len(prs.slides)}", fontsize=11, color="#6b6b6b", va="top")
            ax.text(0.0, 0.88, body, fontsize=12, va="top", linespacing=1.35)
            pdf.savefig(fig, bbox_inches="tight")
            plt.close(fig)

    return True, "fallback_text"


def convert(input_pptx: Path, output_pdf: Path, strategy: str, require_visual: bool) -> tuple[bool, str, list[str]]:
    attempts: list[str] = []
    if strategy == "soffice":
        ok, msg = try_soffice(input_pptx, output_pdf)
        attempts.append(f"soffice:{'ok' if ok else 'fail'}:{msg}")
        return ok, msg, attempts
    if strategy == "keynote":
        ok, msg = try_keynote(input_pptx, output_pdf)
        attempts.append(f"keynote:{'ok' if ok else 'fail'}:{msg}")
        return ok, msg, attempts
    if strategy == "fallback":
        if require_visual:
            return False, "require_visual enabled but strategy=fallback", attempts
        ok, msg = fallback_text_pdf(input_pptx, output_pdf)
        attempts.append(f"fallback:{'ok' if ok else 'fail'}:{msg}")
        return ok, msg, attempts

    ok, mode = try_soffice(input_pptx, output_pdf)
    attempts.append(f"soffice:{'ok' if ok else 'fail'}:{mode}")
    if ok:
        return ok, mode, attempts
    ok, mode2 = try_keynote(input_pptx, output_pdf)
    attempts.append(f"keynote:{'ok' if ok else 'fail'}:{mode2}")
    if ok:
        return ok, mode2, attempts
    if require_visual:
        return False, f"visual conversion failed after soffice({mode}) and keynote({mode2})", attempts
    ok, mode3 = fallback_text_pdf(input_pptx, output_pdf)
    attempts.append(f"fallback:{'ok' if ok else 'fail'}:{mode3}")
    if ok:
        return ok, mode3, attempts
    return False, f"auto failed after soffice({mode}), keynote({mode2}), fallback({mode3})", attempts


def write_report(
    report_path: Path,
    input_pptx: Path,
    output_pdf: Path,
    strategy: str,
    require_visual: bool,
    mode: str,
    attempts: list[str],
) -> None:
    slide_count = 0
    if input_pptx.exists():
        try:
            slide_count = len(Presentation(str(input_pptx)).slides)
        except Exception:
            slide_count = 0
    payload = {
        "timestamp_utc": datetime.now(timezone.utc).isoformat(),
        "input_pptx": str(input_pptx.resolve()),
        "output_pdf": str(output_pdf.resolve()),
        "strategy": strategy,
        "require_visual": require_visual,
        "mode": mode,
        "slide_count": slide_count,
        "output_exists": output_pdf.exists(),
        "output_size_bytes": output_pdf.stat().st_size if output_pdf.exists() else 0,
        "attempts": attempts,
    }
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_text(json.dumps(payload, ensure_ascii=True, indent=2) + "\n", encoding="utf-8")


def main() -> None:
    args = parse_args()
    if not args.input_pptx.exists():
        raise SystemExit(f"Missing input: {args.input_pptx}")
    args.output_pdf.parent.mkdir(parents=True, exist_ok=True)

    ok, mode, attempts = convert(args.input_pptx, args.output_pdf, args.strategy, args.require_visual)
    if not ok:
        if args.report_json is not None:
            write_report(
                args.report_json,
                args.input_pptx,
                args.output_pdf,
                args.strategy,
                args.require_visual,
                f"error:{mode}",
                attempts,
            )
            print(f"Wrote: {args.report_json}")
        raise SystemExit(f"Conversion failed: {mode}")

    if args.report_json is not None:
        write_report(
            args.report_json,
            args.input_pptx,
            args.output_pdf,
            args.strategy,
            args.require_visual,
            mode,
            attempts,
        )
        print(f"Wrote: {args.report_json}")

    print(f"Wrote: {args.output_pdf}")
    print(f"Mode: {mode}")


if __name__ == "__main__":
    main()
