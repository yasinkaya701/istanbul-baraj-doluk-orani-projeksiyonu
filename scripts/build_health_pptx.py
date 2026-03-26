#!/usr/bin/env python3
"""Generate an upgraded 6-slide PowerPoint deck from health-impact outputs.

Enhancements in this version:
- Optional logo support
- Embedded presenter notes in each slide
- Decision matrix and data-coverage cues
- Clearer action framing for 30-60-90 days
"""

from __future__ import annotations

import argparse
from datetime import date
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor(17, 52, 82)
ORANGE = RGBColor(229, 114, 0)
GREEN = RGBColor(36, 133, 74)
RED = RGBColor(184, 47, 47)
GRAY = RGBColor(90, 90, 90)
LIGHT_GRAY = RGBColor(242, 245, 248)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(20, 20, 20)


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Build improved PPTX deck for health-impact summary")
    p.add_argument("--root-dir", type=Path, default=Path("output/health_impact"))
    p.add_argument("--fig-dir", type=Path, default=Path("output/health_impact/figures"))
    p.add_argument(
        "--model-summary-csv",
        type=Path,
        default=None,
        help="Optional model summary CSV path. If omitted, auto-selects the newest available model_comparison_summary*.csv.",
    )
    p.add_argument("--date-label", type=str, default=str(date.today()))
    p.add_argument(
        "--mode",
        choices=["executive", "detailed", "board"],
        default="executive",
        help="executive=6 slides, detailed=10 slides, board=13 slides",
    )
    p.add_argument(
        "--executive-note-md",
        type=Path,
        default=None,
        help="Optional one-liner executive note markdown. If omitted, auto-selects the newest matching file.",
    )
    p.add_argument("--brand-title", type=str, default="Iklim ve Saglik Riski")
    p.add_argument("--logo-path", type=Path, default=None, help="Optional logo image path")
    p.add_argument("--output-path", type=Path, default=None)
    return p.parse_args()


def resolve_model_summary_path(root: Path, override_path: Path | None) -> Path:
    if override_path is not None:
        return override_path
    for p in [
        root / "model_comparison_summary_stable_calibrated.csv",
        root / "model_comparison_summary_duzenlenmis_run.csv",
    ]:
        if p.exists():
            return p
    return root / "model_comparison_summary.csv"


def load_data(root: Path, model_summary_csv: Path | None) -> tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame]:
    model_path = resolve_model_summary_path(root, model_summary_csv)
    model = pd.read_csv(model_path)
    disease = pd.read_csv(root / "health_all_disease_summary.csv")
    matrix = pd.read_csv(root / "health_all_disease_etiology_matrix.csv")
    return model, disease, matrix


def pick_model_row(model: pd.DataFrame, model_name: str) -> pd.Series:
    name = model_name.lower().strip()
    model_labels = model["model"].astype(str).str.lower()
    d = model[model_labels == name]
    if d.empty:
        d = model[model_labels.str.contains(name, regex=False)]
    if d.empty:
        raise ValueError(f"Model not found in model_comparison_summary.csv: {model_name}")
    return d.iloc[0]


def _find_latest(root: Path, pattern: str) -> Path | None:
    matches = sorted(root.glob(pattern))
    if not matches:
        return None
    return matches[-1]


def pick_literature_figure(root: Path, fig_dir: Path, date_label: str) -> Path:
    dated = root / f"literatur_uyum_trafik_isigi_{date_label}.png"
    if dated.exists():
        return dated
    latest = _find_latest(root, "literatur_uyum_trafik_isigi_*.png")
    if latest is not None:
        return latest
    return fig_dir / "fig04_literature_alignment.png"


def load_traffic_light_stats(root: Path, date_label: str) -> dict[str, dict[str, float]]:
    dated = root / f"literatur_uyum_trafik_isigi_{date_label}.csv"
    csv_path = dated if dated.exists() else _find_latest(root, "literatur_uyum_trafik_isigi_*.csv")
    if csv_path is None:
        return {}
    try:
        df = pd.read_csv(csv_path)
    except Exception:
        return {}
    required = {"model", "uyum_etiketi", "oran_yuzde"}
    if not required.issubset(df.columns):
        return {}

    out: dict[str, dict[str, float]] = {}
    model_pairs = [("quant", "quant"), ("strong", "strong")]
    labels = ["yesil_uyumlu", "sari_sinirda", "kirmizi_celiski_riski"]
    for key, needle in model_pairs:
        d = df[df["model"].astype(str).str.lower().str.contains(needle, regex=False)]
        if d.empty:
            continue
        stats = {lbl: 0.0 for lbl in labels}
        for _, row in d.iterrows():
            lbl = str(row["uyum_etiketi"])
            if lbl in stats:
                stats[lbl] = float(row["oran_yuzde"])
        out[key] = stats
    return out


def pick_action_timeline_figure(root: Path, date_label: str) -> Path | None:
    dated = root / f"halk_dili_eylem_takvimi_{date_label}.png"
    if dated.exists():
        return dated
    latest = _find_latest(root, "halk_dili_eylem_takvimi_*.png")
    if latest is not None:
        return latest
    fallback = root / f"halk_dili_oncelik_matrisi_{date_label}.png"
    if fallback.exists():
        return fallback
    fallback_latest = _find_latest(root, "halk_dili_oncelik_matrisi_*.png")
    return fallback_latest


def _pick_csv(root: Path, dated_name: str, glob_pattern: str) -> Path | None:
    dated = root / dated_name
    if dated.exists():
        return dated
    return _find_latest(root, glob_pattern)


def load_action_priority_counts(root: Path, date_label: str) -> dict[str, int]:
    path = _pick_csv(root, f"halk_dili_oncelik_matrisi_{date_label}.csv", "halk_dili_oncelik_matrisi_*.csv")
    if path is None:
        return {}
    try:
        df = pd.read_csv(path)
    except Exception:
        return {}
    if "aksiyon_onceligi" not in df.columns:
        return {}
    vc = df["aksiyon_onceligi"].astype(str).value_counts().to_dict()
    return {
        "acil_eylem": int(vc.get("acil_eylem", 0)),
        "hedefli_onlem": int(vc.get("hedefli_onlem", 0)),
        "yakindan_izlem": int(vc.get("yakindan_izlem", 0)),
        "rutin_izlem": int(vc.get("rutin_izlem", 0)),
    }


def load_action_plan_highlights(root: Path, date_label: str, n: int = 3) -> list[str]:
    path = _pick_csv(root, f"halk_dili_eylem_plani_{date_label}.csv", "halk_dili_eylem_plani_*.csv")
    if path is None:
        return []
    try:
        df = pd.read_csv(path)
    except Exception:
        return []
    required = {"disease_group_tr", "aksiyon_onceligi", "baslama_suresi", "sorumlu_birim"}
    if not required.issubset(df.columns):
        return []
    d = df.head(n).copy()
    lines: list[str] = []
    for _, row in d.iterrows():
        lines.append(
            f"- {row['disease_group_tr']}: {row['aksiyon_onceligi']} | baslama {row['baslama_suresi']} | {row['sorumlu_birim']}"
        )
    return lines


def pick_weekly_workload_figure(root: Path, date_label: str) -> Path | None:
    dated = root / f"halk_dili_haftalik_is_yuku_{date_label}.png"
    if dated.exists():
        return dated
    return _find_latest(root, "halk_dili_haftalik_is_yuku_*.png")


def pick_kpi_alarm_figure(root: Path, date_label: str) -> Path | None:
    dated = root / f"halk_dili_kpi_alarm_panosu_{date_label}.png"
    if dated.exists():
        return dated
    return _find_latest(root, "halk_dili_kpi_alarm_panosu_*.png")


def pick_model_stability_figure(root: Path, date_label: str) -> Path | None:
    dated = root / f"model_stability_dashboard_{date_label}.png"
    if dated.exists():
        return dated
    return _find_latest(root, "model_stability_dashboard_*.png")


def pick_integrated_findings_figure(root: Path, date_label: str) -> Path | None:
    dated = root / f"halk_dili_entegre_bulgular_panosu_{date_label}.png"
    if dated.exists():
        return dated
    return _find_latest(root, "halk_dili_entegre_bulgular_panosu_*.png")


def pick_intervention_scenario_figure(root: Path, date_label: str) -> Path | None:
    dated = root / f"halk_dili_mudahale_senaryolari_{date_label}.png"
    if dated.exists():
        return dated
    return _find_latest(root, "halk_dili_mudahale_senaryolari_*.png")


def pick_operational_readiness_figure(root: Path, date_label: str) -> Path | None:
    dated = root / f"halk_dili_operasyon_hazirlik_panosu_{date_label}.png"
    if dated.exists():
        return dated
    return _find_latest(root, "halk_dili_operasyon_hazirlik_panosu_*.png")


def load_weekly_workload_lines(root: Path, date_label: str, n: int = 3) -> list[str]:
    path = _pick_csv(root, f"halk_dili_haftalik_is_plani_{date_label}.csv", "halk_dili_haftalik_is_plani_*.csv")
    if path is None:
        return []
    try:
        df = pd.read_csv(path)
    except Exception:
        return []
    required = {"hafta", "toplam_aktif_gorev", "acil_eylem", "hedefli_onlem"}
    if not required.issubset(df.columns):
        return []
    out: list[str] = []
    for _, r in df.head(n).iterrows():
        out.append(
            f"- {r['hafta']}: aktif={int(r['toplam_aktif_gorev'])}, "
            f"acil={int(r['acil_eylem'])}, hedefli={int(r['hedefli_onlem'])}"
        )
    return out


def load_intervention_scenario_lines(root: Path, date_label: str, n: int = 2) -> list[str]:
    path = _pick_csv(root, f"halk_dili_mudahale_senaryolari_{date_label}.csv", "halk_dili_mudahale_senaryolari_*.csv")
    if path is None:
        return []
    try:
        df = pd.read_csv(path)
    except Exception:
        return []
    required = {"senaryo", "beklenen_risk_azalimi_yuzde", "kritik_grup_sayisi_kalan"}
    if not required.issubset(df.columns):
        return []
    d = df.sort_values("beklenen_risk_azalimi_yuzde", ascending=False).head(n)
    lines: list[str] = []
    for _, r in d.iterrows():
        lines.append(
            f"- Senaryo {r['senaryo']}: beklenen risk azalimi %{float(r['beklenen_risk_azalimi_yuzde']):.1f}, "
            f"kalan kritik grup {int(r['kritik_grup_sayisi_kalan'])}"
        )
    return lines


def load_operational_readiness_lines(root: Path, date_label: str, n: int = 2) -> list[str]:
    path = _pick_csv(root, f"halk_dili_operasyon_hazirlik_matrisi_{date_label}.csv", "halk_dili_operasyon_hazirlik_matrisi_*.csv")
    if path is None:
        return []
    try:
        df = pd.read_csv(path)
    except Exception:
        return []
    required = {"sorumlu_birim", "hazirlik_baski_puani", "toplam_gorev", "acil_eylem"}
    if not required.issubset(df.columns):
        return []
    d = df.sort_values("hazirlik_baski_puani", ascending=False).head(n)
    lines: list[str] = []
    for _, r in d.iterrows():
        lines.append(
            f"- Birim {r['sorumlu_birim']}: baski {float(r['hazirlik_baski_puani']):.2f}, "
            f"gorev {int(r['toplam_gorev'])}, acil {int(r['acil_eylem'])}"
        )
    return lines


def resolve_executive_note_path(root: Path, date_label: str, override: Path | None) -> Path | None:
    if override is not None:
        return override
    dated = root / f"yonetici_tek_cumle_ozet_{date_label}.md"
    if dated.exists():
        return dated
    latest = _find_latest(root, "yonetici_tek_cumle_ozet_*.md")
    return latest


def load_one_liner_from_md(path: Path | None) -> str | None:
    if path is None or not path.exists():
        return None
    try:
        for line in path.read_text(encoding="utf-8").splitlines():
            s = line.strip()
            if s.startswith("- "):
                return s[2:].strip()
    except Exception:
        return None
    return None


def fallback_one_liner(strong: pd.Series, quant: pd.Series) -> str:
    strong_pct = (float(strong["future_rr_mean"]) - 1.0) * 100.0
    quant_pct = (float(quant["future_rr_mean"]) - 1.0) * 100.0
    return (
        f"2026-2035 projeksiyonunda strong modelde isi-kaynakli risk +%{strong_pct:.1f} "
        f"(RR {float(strong['future_rr_mean']):.3f}), quant modelde +%{quant_pct:.1f} "
        f"(RR {float(quant['future_rr_mean']):.3f}); kapasite planini strong, rutin izlemeyi quant tabaninda yonetin."
    )


def top_disease_lines(disease: pd.DataFrame, model_name: str, n: int = 5) -> list[str]:
    d = disease[disease["model"].str.lower() == model_name.lower()].copy()
    d = d.sort_values("direct_signal_score", ascending=False).head(n)
    return [
        f"{row['disease_group_tr']}: {float(row['direct_signal_score']):.2f} ({row['direct_signal_level']})"
        for _, row in d.iterrows()
    ]


def risk_level(delta_rr: float) -> tuple[str, RGBColor]:
    if delta_rr >= 0.03:
        return "Yuksek", RED
    if delta_rr >= 0.01:
        return "Orta", ORANGE
    return "Dusuk", GREEN


def add_note(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text


def add_shell(prs: Presentation, title: str, subtitle: str, date_label: str, logo_path: Path | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.95), Inches(13.333), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.45), Inches(0.13), Inches(9.2), Inches(0.7))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if logo_path and logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(11.2), Inches(0.14), height=Inches(0.62))

    sub = slide.shapes.add_textbox(Inches(0.45), Inches(1.08), Inches(10.5), Inches(0.45))
    stf = sub.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(14)
    sp.font.color.rgb = GRAY

    foot = slide.shapes.add_textbox(Inches(0.45), Inches(7.12), Inches(12.4), Inches(0.22))
    ftf = foot.text_frame
    ftf.clear()
    fp = ftf.paragraphs[0]
    fp.text = f"Kaynak: output/health_impact | Tarih: {date_label}"
    fp.font.size = Pt(9)
    fp.font.color.rgb = GRAY
    fp.alignment = PP_ALIGN.RIGHT

    return slide


def add_kpi_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    value: str,
    sub: str,
    color: RGBColor,
):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.color.rgb = color
    card.line.width = Pt(1.8)

    tx = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.08), Inches(width - 0.3), Inches(0.25))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = GRAY

    vx = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.42), Inches(width - 0.3), Inches(0.5))
    vtf = vx.text_frame
    vtf.clear()
    vp = vtf.paragraphs[0]
    vp.text = value
    vp.font.size = Pt(25)
    vp.font.bold = True
    vp.font.color.rgb = color

    sx = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.94), Inches(width - 0.3), Inches(0.35))
    stf = sx.text_frame
    stf.clear()
    sp = stf.paragraphs[0]
    sp.text = sub
    sp.font.size = Pt(11)
    sp.font.color.rgb = BLACK


def add_bullets(slide, left: float, top: float, width: float, height: float, bullets: list[str], font_size: int = 18):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.clear()
    for i, txt in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.level = 0


def add_figure(slide, fig_path: Path, left: float, top: float, width: float, height: float, caption: str):
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = RGBColor(205, 212, 219)
    frame.line.width = Pt(1.0)

    if fig_path.exists():
        slide.shapes.add_picture(
            str(fig_path),
            Inches(left + 0.08),
            Inches(top + 0.08),
            width=Inches(width - 0.16),
            height=Inches(height - 0.42),
        )
    tx = slide.shapes.add_textbox(Inches(left + 0.12), Inches(top + height - 0.28), Inches(width - 0.25), Inches(0.2))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = caption
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY


def add_decision_matrix(slide, left: float, top: float, width: float, height: float, strong_lvl: str, quant_lvl: str):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.color.rgb = NAVY
    box.line.width = Pt(1.2)

    title = slide.shapes.add_textbox(Inches(left + 0.14), Inches(top + 0.08), Inches(width - 0.25), Inches(0.28))
    ttf = title.text_frame
    ttf.clear()
    tp = ttf.paragraphs[0]
    tp.text = "Karar Matrisi"
    tp.font.size = Pt(12)
    tp.font.bold = True
    tp.font.color.rgb = NAVY

    bullets = [
        f"- Strong risk: {strong_lvl}",
        f"- Quant risk: {quant_lvl}",
        "- Operasyon tercihi: strong-senaryo hazirlik, quant-senaryo izleme",
    ]
    add_bullets(slide, left + 0.16, top + 0.43, width - 0.3, height - 0.5, bullets, font_size=10)


def build_pptx(args: argparse.Namespace) -> Path:
    model, disease, matrix = load_data(args.root_dir, args.model_summary_csv)
    strong = pick_model_row(model, "strong")
    quant = pick_model_row(model, "quant")

    strong_top = top_disease_lines(disease, "strong", 5)
    quant_top = top_disease_lines(disease, "quant", 5)
    strong_level, strong_color = risk_level(float(strong["delta_rr_mean"]))
    quant_level, quant_color = risk_level(float(quant["delta_rr_mean"]))
    direct_share = float(matrix["quantifiable_with_current_data"].mean())

    if args.output_path is not None:
        out_path = args.output_path
    elif args.mode == "detailed":
        out_path = args.root_dir / f"sunum_10_slayt_{args.date_label}_v4_detailed.pptx"
    elif args.mode == "board":
        out_path = args.root_dir / f"sunum_13_slayt_{args.date_label}_v6_board.pptx"
    else:
        out_path = args.root_dir / f"sunum_6_slayt_{args.date_label}_v3.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    s1 = add_shell(
        prs,
        f"{args.brand_title} | Yonetici Sunumu",
        "Ana mesaj: strong senaryoda risk artisi belirgin, quant senaryoda sinirli.",
        args.date_label,
        args.logo_path,
    )
    add_kpi_card(s1, 0.7, 1.9, 3.9, 1.45, "Strong Gelecek RR", f"{float(strong['future_rr_mean']):.3f}", "Isi etkili goreli risk", strong_color)
    add_kpi_card(s1, 4.9, 1.9, 3.9, 1.45, "Quant Gelecek RR", f"{float(quant['future_rr_mean']):.3f}", "Isi etkili goreli risk", quant_color)
    add_kpi_card(s1, 9.1, 1.9, 3.5, 1.45, "Risk Seviye", f"S:{strong_level} / Q:{quant_level}", "Delta RR bazli", ORANGE)
    add_bullets(
        s1,
        0.75,
        3.75,
        11.8,
        2.8,
        [
            f"- Esik asimi ay payi: strong {float(strong['future_threshold_exceed_share']):.1%} | quant {float(quant['future_threshold_exceed_share']):.1%}",
            f"- OOD artis payi: strong {float(strong['delta_ood_share']):.1%} | quant {float(quant['delta_ood_share']):.1%}",
            "- Bu projeksiyonda genel kurulasma sinyali yok; bagil nem artisi goruluyor.",
        ],
        font_size=17,
    )
    add_note(
        s1,
        "Bu slaytta karar mesajini net ver: strong yuksek baski, quant dusuk baski. "
        "Kritik KPI'lari bir dakikada gec.",
    )

    # Slide 2
    s2 = add_shell(prs, "Sayisal Durum Panosu", "Baz: 1991-2020 | Gelecek: 2026-2035", args.date_label, args.logo_path)
    add_kpi_card(s2, 0.7, 1.7, 4.0, 1.4, "Strong Delta RR", f"{float(strong['delta_rr_mean']):+.4f}", "Daha yuksek artis", strong_color)
    add_kpi_card(s2, 4.95, 1.7, 4.0, 1.4, "Quant Delta RR", f"{float(quant['delta_rr_mean']):+.4f}", "Sinirli artis", quant_color)
    add_kpi_card(s2, 9.2, 1.7, 3.5, 1.4, "AF Farki", f"{float(strong['future_af_mean']):.3f} vs {float(quant['future_af_mean']):.3f}", "Strong vs Quant", ORANGE)
    add_kpi_card(s2, 0.7, 3.35, 4.0, 1.4, "HI Delta (C)", f"{float(strong['delta_hi_mean_c']):.2f} / {float(quant['delta_hi_mean_c']):.2f}", "Strong / Quant", RED)
    add_kpi_card(s2, 4.95, 3.35, 4.0, 1.4, "Threshold Exceed", f"{float(strong['future_threshold_exceed_share']):.1%} / {float(quant['future_threshold_exceed_share']):.1%}", "Strong / Quant", ORANGE)
    add_kpi_card(s2, 9.2, 3.35, 3.5, 1.4, "OOD Payi", f"{float(strong['delta_ood_share']):.1%} / {float(quant['delta_ood_share']):.1%}", "Strong / Quant", RED)
    add_decision_matrix(s2, 8.95, 5.03, 3.75, 1.8, strong_level, quant_level)
    add_bullets(
        s2,
        0.75,
        5.05,
        7.9,
        1.8,
        [
            "- Planlama notu: guvenlik payi gerekiyorsa strong bazli kapasite planla.",
            "- Kaynak kisitinda quant tabanli minimum plan + tetikleyici esik kullan.",
        ],
        font_size=14,
    )
    add_note(
        s2,
        "Sayisal kartlari karsilastirarak anlat. Decision matrix bolumunde "
        "hangi senaryoda hangi operasyon stratejisini sececegini acikla.",
    )

    # Slide 3
    s3 = add_shell(prs, "Model Ayrisimi", "Risk farkinin gorunur ozeti", args.date_label, args.logo_path)
    add_bullets(
        s3,
        0.75,
        1.7,
        4.45,
        4.9,
        [
            "- Iki model farkli risk seviyesi uretiyor.",
            "- Strong: yuksek termal baski ve daha fazla esik asimi.",
            "- Quant: artis zayif ancak sifir degil.",
            "- Duyarlilikte strong 162/162 pozitif delta RR.",
        ],
        font_size=17,
    )
    add_figure(
        s3,
        args.fig_dir / "fig01_model_overview.png",
        left=5.35,
        top=1.55,
        width=7.55,
        height=5.5,
        caption="Model ozet gorunumu",
    )
    add_note(
        s3,
        "Grafikten tek mesaj: belirsizlik olsa da strong senaryoda artis kalici. "
        "Quant model alt sinir davranisini temsil ediyor.",
    )

    # Slide 4
    s4 = add_shell(prs, "Hastalik Gruplari Etkisi", "Dogrudan iklim sinyaline gore onceliklendirme", args.date_label, args.logo_path)
    add_bullets(
        s4,
        0.75,
        1.65,
        5.3,
        4.9,
        ["Strong ilk 5:"] + [f"- {x}" for x in strong_top] + ["", "Quant ilk 5:"] + [f"- {x}" for x in quant_top],
        font_size=13,
    )
    add_figure(
        s4,
        args.fig_dir / "fig06_all_disease_direct_scores.png",
        left=6.1,
        top=1.55,
        width=6.8,
        height=5.5,
        caption="Dogrudan risk sinyal skorlari",
    )
    add_note(
        s4,
        "Oncelikli alanlar: strong modelde genel isi yuk, kardiyovaskuler, bobrek, anne-bebek, solunum. "
        "Kaynak dagitiminda bu siralamayi kullan.",
    )

    # Slide 5
    s5 = add_shell(prs, "Etiyoloji Kapsami ve Sinirlar", "Neyi olcebiliyoruz, neleri nitel izliyoruz", args.date_label, args.logo_path)
    add_kpi_card(
        s5,
        0.75,
        1.72,
        5.0,
        1.2,
        "Dogrudan Olculebilir Kapsam",
        f"{direct_share:.1%}",
        "Mevcut veri seti (isi+nem) ile",
        NAVY,
    )
    add_bullets(
        s5,
        0.75,
        3.05,
        5.0,
        3.4,
        [
            "- Dogrudan: isi ve nem etkisi.",
            "- Nitel: UV, PM2.5/O3, polen, vektor, su-gida, afet kesintisi.",
            "- Cilt kanseri ana etiyoloji: UV.",
            "- UV katmani olmadan sayisal cilt kanseri artis tahmini verilmez.",
        ],
        font_size=14,
    )
    add_figure(
        s5,
        args.fig_dir / "fig07_all_disease_etiology_coverage.png",
        left=5.9,
        top=1.55,
        width=7.0,
        height=5.5,
        caption="1=dogrudan hesaplanabilir | 0=nitel",
    )
    add_note(
        s5,
        "Bu slayt metodolojik sinirlari aciklar. Ozel vurgu: cilt kanseri icin UV gerekir; "
        "mevcut model bu yolu nitel izliyor.",
    )

    # Slide 6
    s6 = add_shell(prs, "90 Gunluk Uygulama Yol Haritasi", "Uygulanabilir ve olculebilir adimlar", args.date_label, args.logo_path)
    box_specs = [
        (0.8, 1.8, 4.0, 4.8, "Ilk 30 Gun", ORANGE, ["- Veri entegrasyonu", "- Gunluk saglik sonlanim tablosu", "- Erken uyari kriter taslagi"]),
        (4.95, 1.8, 4.0, 4.8, "30-60 Gun", NAVY, ["- DLNM kalibrasyonu", "- UV + PM2.5 + O3 + polen katmani", "- Ilce bazli risk haritalama"]),
        (9.1, 1.8, 3.4, 4.8, "60-90 Gun", GREEN, ["- Operasyonel panel yayini", "- Kirmizi grup protokolleri", "- Aylik KPI denetimi"]),
    ]
    for left, top, width, height, title, color, items in box_specs:
        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = color
        card.line.width = Pt(2.0)

        tx = s6.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.12), Inches(width - 0.25), Inches(0.35))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        add_bullets(s6, left + 0.15, top + 0.58, width - 0.28, height - 0.8, items, font_size=13)

    add_note(
        s6,
        "Yol haritasini kaynak ve sorumlulukla bagla. 90. gun sonunda hedef: "
        "calisan erken uyari paneli + kalibre model + koruma protokolleri.",
    )

    if args.mode in {"detailed", "board"}:
        hum_delta_strong = float(disease[disease["model"].str.lower() == "strong"]["delta_mean_humidity_pct"].iloc[0])
        hum_delta_quant = float(disease[disease["model"].str.lower() == "quant"]["delta_mean_humidity_pct"].iloc[0])
        traffic = load_traffic_light_stats(args.root_dir, args.date_label)
        lit_fig = pick_literature_figure(args.root_dir, args.fig_dir, args.date_label)

        if "quant" in traffic:
            quant_line = (
                f"- Quant trafik: yesil %{traffic['quant']['yesil_uyumlu']:.1f}, "
                f"sari %{traffic['quant']['sari_sinirda']:.1f}, "
                f"kirmizi %{traffic['quant']['kirmizi_celiski_riski']:.1f}."
            )
        else:
            quant_line = "- Quant trafik ozeti bulunamadi (CSV eksik)."

        if "strong" in traffic:
            strong_line = (
                f"- Strong trafik: yesil %{traffic['strong']['yesil_uyumlu']:.1f}, "
                f"sari %{traffic['strong']['sari_sinirda']:.1f}, "
                f"kirmizi %{traffic['strong']['kirmizi_celiski_riski']:.1f}."
            )
        else:
            strong_line = "- Strong trafik ozeti bulunamadi (CSV eksik)."

        # Slide 7
        s7 = add_shell(prs, "Senaryo Dagilimi ve Belirsizlik", "Duyarlilik analizi ozeti", args.date_label, args.logo_path)
        add_bullets(
            s7,
            0.75,
            1.7,
            4.9,
            4.9,
            [
                "- Senaryo dagilimlari model secimine hassas.",
                "- Strong tarafinda risk araligi daha genis.",
                "- Quant tarafinda aralik daha dar ama sifir degil.",
                "- Politika kararinda belirsizlik bandi dikkate alinmali.",
            ],
            font_size=16,
        )
        add_figure(
            s7,
            args.fig_dir / "fig03_scenario_ranges.png",
            left=5.8,
            top=1.55,
            width=7.0,
            height=5.5,
            caption="Future RR senaryo araliklari",
        )
        add_note(s7, "Belirsizlik slaydi: tek sayiya degil, araliklara gore karar verilmesi gerektigini vurgula.")

        # Slide 8
        s8 = add_shell(prs, "Literatur Uyum ve Kanit", "Model sinyali literaturle uyumlu mu?", args.date_label, args.logo_path)
        add_bullets(
            s8,
            0.75,
            1.7,
            5.1,
            4.9,
            [
                "- Isi-saglik etkisi literaturde tutarli sekilde gosterilmistir.",
                strong_line,
                quant_line,
                "- Proxy model siniri: UV/PM2.5/polen katmani ayrica gerekir.",
            ],
            font_size=15,
        )
        add_figure(
            s8,
            lit_fig,
            left=5.95,
            top=1.55,
            width=6.85,
            height=5.5,
            caption="Literatur uyum trafik isigi",
        )
        add_note(s8, "Bu slayt guven olusturur: modelin ana yonu literaturle celismiyor mesajini ver.")

        # Slide 9
        s9 = add_shell(prs, "Cilt Kanseri ve Kuruluk Ozel Notu", "Sik sorulan iki soru", args.date_label, args.logo_path)
        add_kpi_card(s9, 0.75, 1.75, 5.0, 1.3, "Nem Delta (Strong/Quant)", f"{hum_delta_strong:+.2f} / {hum_delta_quant:+.2f}", "Bagil nem puan degisimi", NAVY)
        add_bullets(
            s9,
            0.75,
            3.2,
            5.1,
            3.6,
            [
                "- Bu projeksiyonda genel kurulasma sinyali yok.",
                "- Cilt kanseri ana etiyoloji UV maruziyetidir.",
                "- UV katmani olmadan sayisal cilt kanseri artisi verilmez.",
                "- Mevcut cikti: nitel risk baskisi + veri gereksinimi uyarisi.",
            ],
            font_size=15,
        )
        add_figure(
            s9,
            args.fig_dir / "fig07_all_disease_etiology_coverage.png",
            left=5.9,
            top=1.55,
            width=6.9,
            height=5.5,
            caption="UV gibi etiyolojiler icin ek veri gereksinimi",
        )
        add_note(s9, "Soru-Cevapta en cok gelecek konu: cilt kanseri ve kuruluk. Bu slaydi referans ver.")

        # Slide 10
        s10 = add_shell(prs, "Kapanis ve Karar", "Yonetici icin net secim ciktisi", args.date_label, args.logo_path)
        add_bullets(
            s10,
            0.75,
            1.8,
            12.0,
            4.8,
            [
                "- Karar 1: Operasyonel kapasiteyi strong-senaryo guvenlik payiyla planla.",
                "- Karar 2: Quant-senaryoyu alt-sinir kontrolu olarak aylik izle.",
                "- Karar 3: 90 gun icinde UV + PM2.5/O3 + polen katmanlarini entegre et.",
                "- Karar 4: Erken uyari panelini ilce/hafta bazli canliya al.",
                "- Karar 5: Kirmizi gruplar icin hedefli koruma protokollerini devreye al.",
            ],
            font_size=17,
        )
        add_note(s10, "Kapanista 5 karar maddesini okuyup toplantiyi aksiyonla kapat.")

    if args.mode == "board":
        exec_note_path = resolve_executive_note_path(args.root_dir, args.date_label, args.executive_note_md)
        one_liner = load_one_liner_from_md(exec_note_path) or fallback_one_liner(strong, quant)
        traffic = load_traffic_light_stats(args.root_dir, args.date_label)
        action_counts = load_action_priority_counts(args.root_dir, args.date_label)
        action_highlights = load_action_plan_highlights(args.root_dir, args.date_label, n=3)
        action_fig = pick_action_timeline_figure(args.root_dir, args.date_label)
        integrated_fig = pick_integrated_findings_figure(args.root_dir, args.date_label)
        scenario_fig = pick_intervention_scenario_figure(args.root_dir, args.date_label)
        ops_fig = pick_operational_readiness_figure(args.root_dir, args.date_label)
        weekly_fig = pick_weekly_workload_figure(args.root_dir, args.date_label)
        kpi_alarm_fig = pick_kpi_alarm_figure(args.root_dir, args.date_label)
        stability_fig = pick_model_stability_figure(args.root_dir, args.date_label)
        weekly_lines = load_weekly_workload_lines(args.root_dir, args.date_label, n=3)
        scenario_lines = load_intervention_scenario_lines(args.root_dir, args.date_label, n=2)
        ops_lines = load_operational_readiness_lines(args.root_dir, args.date_label, n=2)
        quant_tr = traffic.get("quant", {})
        strong_tr = traffic.get("strong", {})
        strong_line = (
            f"- Strong literatur trafik: yesil %{float(strong_tr.get('yesil_uyumlu', 0.0)):.1f}, "
            f"sari %{float(strong_tr.get('sari_sinirda', 0.0)):.1f}, "
            f"kirmizi %{float(strong_tr.get('kirmizi_celiski_riski', 0.0)):.1f}."
        )
        quant_line = (
            f"- Quant literatur trafik: yesil %{float(quant_tr.get('yesil_uyumlu', 0.0)):.1f}, "
            f"sari %{float(quant_tr.get('sari_sinirda', 0.0)):.1f}, "
            f"kirmizi %{float(quant_tr.get('kirmizi_celiski_riski', 0.0)):.1f}."
        )

        # Slide 11: Governance and ownership + 90-day execution timeline
        s11 = add_shell(prs, "Yonetişim ve Sorumluluk Matrisi", "Kim neyi ne zaman yapacak?", args.date_label, args.logo_path)
        action_summary = [
            f"- Aksiyon dagilimi: acil={int(action_counts.get('acil_eylem', 0))}, "
            f"hedefli={int(action_counts.get('hedefli_onlem', 0))}, "
            f"yakindan={int(action_counts.get('yakindan_izlem', 0))}, "
            f"rutin={int(action_counts.get('rutin_izlem', 0))}.",
        ]
        if action_highlights:
            action_summary.extend(action_highlights)
        else:
            action_summary.append("- Eylem plani CSV bulunamadi; genel yonetisim maddeleri kullaniliyor.")
        if ops_lines:
            action_summary.extend(ops_lines)

        add_bullets(
            s11,
            0.75,
            1.8,
            5.45,
            4.8,
            action_summary
            + [
                "- Veri ve model yonetimi: Teknik ekip (haftalik kalite raporu).",
                "- Erken uyari ve saha koordinasyonu: Kamu sagligi / afet birimi.",
                "- KPI denetimi: Aylik yonetici komite toplantisi.",
            ],
            font_size=12,
        )
        governance_fig = ops_fig or integrated_fig or scenario_fig or action_fig
        if governance_fig is not None:
            add_figure(
                s11,
                governance_fig,
                left=6.35,
                top=1.55,
                width=6.55,
                height=5.5,
                caption="Operasyon hazirlik panosu (varsa) / entegre bulgular / eylem takvimi",
            )
        add_note(
            s11,
            "Bu slaytta eylem dagilimini ve ilk 90 gun takvimini goster. "
            "Acil-hedefli-rutin dagilimini sorumlu birimlerle baglayarak anlat.",
        )

        # Slide 12: Q&A + weekly monitoring
        s12 = add_shell(prs, "SSS ve Haftalik Izlem", "Toplanti sonu itiraz yonetimi + is yuku takibi", args.date_label, args.logo_path)
        weekly_block = weekly_lines if weekly_lines else ["- Haftalik is yuku tablosu bulunamadi."]
        add_bullets(
            s12,
            0.75,
            1.8,
            5.5,
            4.8,
            [
                "- Soru: 'Model neden iki farkli sonuc veriyor?' Yanit: Belirsizlik bandi/varsayim farki.",
                "- Soru: 'Cilt kanseri artis yuzdesi neden yok?' Yanit: UV katmani olmadan sayisal tahmin verilmez.",
                "- Soru: 'Kuruluk var mi?' Yanit: Bu projeksiyonda bagil nem artisi var, genel kurulasma sinyali yok.",
                "- Soru: 'Hangi senaryoya gore aksiyon?' Yanit: Kapasite plani strong, izleme tabani quant.",
                "- Soru: '90 gunde ne bitecek?' Yanit: Erken uyari paneli + coklu maruziyet katmani + KPI dongusu.",
                "- Soru: 'Alarm neye gore calisir?' Yanit: KPI hedeflerinin alt/ust esik sapmalarina gore sari-kirmizi alarm.",
                "- Soru: 'Tahmin neden daha stabil?' Yanit: Duyarlilik dagilimi + OOD sinyali ile kalibrasyon uygulaniyor.",
                "",
                "Haftalik izlem (ilk haftalar):",
            ]
            + weekly_block,
            font_size=12,
        )
        monitor_fig = kpi_alarm_fig or stability_fig or weekly_fig
        if monitor_fig is not None:
            add_figure(
                s12,
                monitor_fig,
                left=6.35,
                top=1.55,
                width=6.55,
                height=5.5,
                caption="KPI alarm / model stabilite / haftalik aktif gorev dagilimi",
            )
        add_note(
            s12,
            "SSS yanitlarini haftalik is yukune bagla. "
            "Hafta 1-2'nin acil/ hedefli dagilimini yonetime net goster.",
        )

        # Slide 13: One-liner decision message
        s13 = add_shell(prs, "Tek Cumle Karar Mesaji", "Toplanti kapanisinda okunacak net mesaj", args.date_label, args.logo_path)
        msg_box = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.8), Inches(11.6), Inches(2.1))
        msg_box.fill.solid()
        msg_box.fill.fore_color.rgb = LIGHT_GRAY
        msg_box.line.color.rgb = NAVY
        msg_box.line.width = Pt(2.0)

        tx = s13.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11.0), Inches(1.5))
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = one_liner
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = BLACK

        scenario_block = scenario_lines if scenario_lines else ["- Mudahale senaryo tablosu bulunamadi."]
        add_bullets(
            s13,
            1.1,
            4.35,
            11.3,
            2.3,
            [
                "- Uygulama kurali: kapasite planlama strong, operasyonel taban izleme quant.",
                strong_line,
                quant_line,
            ]
            + scenario_block
            + [
                "- Teknik not: UV + PM2.5/O3 + polen katmanlari olmadan etiyoloji kapsami sinirlidir.",
            ],
            font_size=14,
        )
        add_note(
            s13,
            "Bu slaydi son 30 saniyede oku. Amac: tek cümlede karar ve sonraki adimi netlestirmek.",
        )

    prs.save(str(out_path))
    return out_path


def main() -> None:
    args = parse_args()
    out = build_pptx(args)
    print(f"Wrote: {out}")


if __name__ == "__main__":
    main()
